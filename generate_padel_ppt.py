#!/usr/bin/env python3
"""
Genera la presentación PPT «Academia Pádel — SaaS de Gestión»
usando la plantilla Ai Automation Business ppt.pptx como base visual.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json, os

# ──────────────────────────────────────────────
# Paths
# ──────────────────────────────────────────────
TEMPLATE = "/docker/data_dockers/Ai Automation Business ppt.pptx"
SCREENSHOTS = "/docker/data_dockers/padel_app/screenshots"
OUTPUT = "/docker/data_dockers/padel_app/Academia_Padel_Presentacion.pptx"
LOGO_PATH = "/docker/data_dockers/_assets/ALTER SIN NOMBRE.png"

# ──────────────────────────────────────────────
# Color palette (from the template style)
# ──────────────────────────────────────────────
PRIMARY    = RGBColor(0x1E, 0x83, 0xEC)  # App primary blue
DARK       = RGBColor(0x0F, 0x17, 0x2A)  # Dark bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF4, 0xF6, 0xF8)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
EMERALD    = RGBColor(0x22, 0xC5, 0x5E)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
RED        = RGBColor(0xEF, 0x44, 0x44)
VIOLET     = RGBColor(0x8B, 0x5C, 0xF6)
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE    = RGBColor(0xE8, 0xEF, 0xF4)
TEAL       = RGBColor(0x0F, 0x8B, 0x8D)
AMBER_ALT  = RGBColor(0xF2, 0xA6, 0x40)

# ──────────────────────────────────────────────
# Data from API (captured values)
# ──────────────────────────────────────────────
DATA = {
    "totalAlumnos": 14,
    "totalClases": 8,
    "recuperaciones": 4,
    "sesionesEstaSemana": 5,
    "profesores": 3,
    "pistas": 3,
    "niveles": ["Iniciación", "Intermedio", "Avanzado", "Competición"],
    "mrr": 826,  # 14 * 59
    "ocupacion": "44%",
    "retencion": "94%",
    "clases": [
        ("Iniciación Mañana", "Lunes 09:00-10:00", "Pista 1", "Carlos", "3/4"),
        ("Competición", "Lunes 19:00-20:30", "Pista 1", "Laura", "2/4"),
        ("Iniciación Tarde", "Martes 18:00-19:00", "Pista 2", "Laura", "3/4"),
        ("Intermedio Mañana", "Miércoles 10:00-11:00", "Pista 1", "Carlos", "1/4"),
        ("Competición", "Miércoles 19:00-20:30", "Pista 1", "Laura", "0/4"),
        ("Intermedio Tarde", "Jueves 19:00-20:00", "Pista 3", "Sergio", "2/4"),
        ("Competición", "Viernes 19:00-20:30", "Pista 1", "Laura", "0/4"),
        ("Avanzado Sábado", "Sábado 10:00-11:30", "Pista 2", "Laura", "3/4"),
    ],
}

# ──────────────────────────────────────────────
# Load template & setup
# ──────────────────────────────────────────────
prs = Presentation(TEMPLATE)
W = prs.slide_width   # 13.33 in
H = prs.slide_height  # 7.5 in

# Delete all existing slides from the template
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        prs.part.drop_rel(rId)
    xml_slides.remove(sldId)

# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, color, radius=0.05):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def text(slide, left, top, width, height, txt, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tf

def bullets(slide, left, top, width, height, items, size=12, color=DARK, spacing=4):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tf

def add_footer(slide):
    rect(slide, 0, H - Inches(0.4), W, Inches(0.4), DARK)
    text(slide, 0, H - Inches(0.38), W, Inches(0.35),
         "Grupo ALTER  ·  Departamento de Informática  ·  Marzo 2026  ·  Confidencial",
         size=8, color=MUTED, align=PP_ALIGN.CENTER)

def accent_line(slide, left, top, width, color=PRIMARY):
    rect(slide, left, top, width, Inches(0.04), color)

def kpi_card(slide, left, top, value, label, accent=PRIMARY):
    w, h = Inches(2.6), Inches(1.6)
    box(slide, left, top, w, h, WHITE)
    accent_line(slide, left + Inches(0.3), top + Inches(0.12), Inches(2.0), accent)
    text(slide, left, top + Inches(0.25), w, Inches(0.7), str(value),
         size=36, color=accent, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left, top + Inches(1.0), w, Inches(0.5), label,
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

def add_screenshot(slide, img_file, left, top, width, height=None):
    path = f"{SCREENSHOTS}/{img_file}"
    if not os.path.exists(path):
        return
    if height:
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        slide.shapes.add_picture(path, left, top, width=width)

def section_tag(slide, left, top, label, color=PRIMARY):
    """Small colored tag/badge."""
    w = Inches(len(label) * 0.09 + 0.4)
    h = Inches(0.3)
    b = box(slide, left, top, w, h, color, radius=0.15)
    tf = b.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"


# ═══════════════════════════════════════════════
# SLIDE 1 — Portada
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])  # blank
add_bg(slide, DARK)

# Top accent bars
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)
rect(slide, 0, Inches(0.05), W, Inches(0.03), AMBER)

# Title
text(slide, Inches(1), Inches(1.6), Inches(11), Inches(1.0),
     "Academia Pádel",
     size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.7),
     "Plataforma SaaS de Gestión para Academias de Pádel",
     size=24, color=PRIMARY, align=PP_ALIGN.CENTER)

# Accent line
accent_line(slide, Inches(4.5), Inches(3.7), Inches(4.3), AMBER)

# Info
text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
     "Grupo ALTER  ·  Departamento de Informática  ·  Marzo 2026",
     size=15, color=MUTED, align=PP_ALIGN.CENTER)

# Badge with stats
text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
     f"🎾  {DATA['totalAlumnos']} Alumnos  ·  {DATA['totalClases']} Clases  ·  {DATA['profesores']} Profesores  ·  {DATA['pistas']} Pistas  ·  3 Planes SaaS",
     size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Footer
text(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.5),
     "Next.js · Express · PostgreSQL · Docker · Tailwind CSS",
     size=11, color=RGBColor(0x4A, 0x5A, 0x6A), align=PP_ALIGN.CENTER)

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 2 — Índice / Agenda
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, LIGHT_BG)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.5),
     "Agenda", size=14, color=PRIMARY, bold=True)
text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
     "Qué veremos", size=36, color=DARK, bold=True)

items = [
    ("01", "Visión General", "Qué es y qué resuelve la plataforma"),
    ("02", "Arquitectura Técnica", "Stack tecnológico y despliegue Docker"),
    ("03", "Modelo SaaS", "Planes Starter, Club y Elite"),
    ("04", "Funcionalidades", "Alumnos, clases, calendario, recuperaciones..."),
    ("05", "Dashboard Elite", "Estadísticas avanzadas y gráficos"),
    ("06", "Datos Reales", "KPIs del sistema en producción"),
    ("07", "Pantallazos", "Capturas de la aplicación real"),
]

y_start = Inches(2.0)
for i, (num, title, desc) in enumerate(items):
    y = y_start + Inches(i * 0.7)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    run.font.bold = True

    text(slide, Inches(2.0), y - Inches(0.02), Inches(4), Inches(0.35),
         title, size=16, color=DARK, bold=True)
    text(slide, Inches(2.0), y + Inches(0.28), Inches(8), Inches(0.3),
         desc, size=11, color=MUTED)

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 3 — Visión General
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, LIGHT_BG)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.5), "VISIÓN GENERAL")
text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
     "¿Qué es Academia Pádel?", size=34, color=DARK, bold=True)

text(slide, Inches(0.8), Inches(1.9), Inches(6), Inches(1.5),
     "Plataforma SaaS completa para la gestión integral de academias de pádel. "
     "Permite gestionar alumnos, clases, horarios, asistencia, recuperaciones, "
     "pagos y comunicaciones desde un único panel web.",
     size=13, color=MUTED)

# Problem → Solution cards
y_cards = Inches(3.5)
card_w = Inches(5.5)
card_h = Inches(2.8)

# Problem card
box(slide, Inches(0.8), y_cards, card_w, card_h, WHITE)
accent_line(slide, Inches(1.1), y_cards + Inches(0.15), Inches(2.0), RED)
text(slide, Inches(1.1), y_cards + Inches(0.3), Inches(5), Inches(0.4),
     "❌  Problema", size=16, color=RED, bold=True)
bullets(slide, Inches(1.1), y_cards + Inches(0.8), Inches(5), Inches(2.0), [
    "Gestión manual con Excel y papel",
    "Sin control de asistencia ni recuperaciones",
    "Dificultad para comunicar cambios a familias",
    "Cero visibilidad sobre ocupación y rentabilidad",
], size=11, color=MUTED)

# Solution card
box(slide, Inches(7.0), y_cards, card_w, card_h, WHITE)
accent_line(slide, Inches(7.3), y_cards + Inches(0.15), Inches(2.0), EMERALD)
text(slide, Inches(7.3), y_cards + Inches(0.3), Inches(5), Inches(0.4),
     "✅  Solución", size=16, color=EMERALD, bold=True)
bullets(slide, Inches(7.3), y_cards + Inches(0.8), Inches(5), Inches(2.0), [
    "Panel web centralizado con datos en tiempo real",
    "Asistencia automática con recuperaciones",
    "Notificaciones email y WhatsApp integradas",
    "Dashboard con KPIs, gráficos y estadísticas",
], size=11, color=MUTED)

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 4 — Arquitectura Técnica
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.5), "ARQUITECTURA")
text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
     "Stack Tecnológico", size=34, color=WHITE, bold=True)

text(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(0.5),
     "Aplicación full-stack moderna con despliegue Docker",
     size=14, color=MUTED)

# Tech stack cards
cards = [
    ("Frontend", "Next.js 15\nTypeScript\nTailwind CSS v4\nLucide Icons", PRIMARY),
    ("Backend", "Express v5\nTypeScript\nPrisma v6 ORM\nREST API", EMERALD),
    ("Base de Datos", "PostgreSQL 16\nMigraciones Prisma\nRelaciones complejas\nBackups automáticos", VIOLET),
    ("Infraestructura", "Docker Compose\n3 contenedores\nProxy Nginx\nSSL/HTTPS", AMBER),
]

card_w = Inches(2.8)
card_h = Inches(3.5)
gap = Inches(0.3)
start_x = Inches(0.6)

for i, (title, items_txt, accent) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    y = Inches(2.6)
    b = box(slide, x, y, card_w, card_h, RGBColor(0x16, 0x1E, 0x30))
    accent_line(slide, x + Inches(0.3), y + Inches(0.15), Inches(2.0), accent)
    text(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.4),
         title, size=16, color=accent, bold=True)
    # Items
    for j, line in enumerate(items_txt.split("\n")):
        text(slide, x + Inches(0.3), y + Inches(0.8) + Inches(j * 0.55), Inches(2.2), Inches(0.4),
             f"▸  {line}", size=11, color=RGBColor(0x9A, 0xB0, 0xBD))

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 5 — Modelo SaaS (3 planes)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, LIGHT_BG)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.5), "MODELO SAAS")
text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
     "Tres planes adaptados a cada club", size=34, color=DARK, bold=True)

plans = [
    ("Starter", "39€/mes", EMERALD, [
        "Gestión de clientes",
        "Clases y horarios",
        "Niveles de jugadores",
        "Control de asistencia",
        "Notificaciones email",
    ]),
    ("Club", "59€/mes", AMBER, [
        "Todo de Starter +",
        "Recuperación de clases",
        "Asignación automática",
        "Control de pagos",
        "Vista diaria del club",
    ]),
    ("Elite", "79€/mes", RED, [
        "Todo de Club +",
        "Automatizaciones avanzadas",
        "Reporting avanzado",
        "Multi-club",
        "Integraciones (API)",
        "Personalización (branding)",
    ]),
]

card_w = Inches(3.8)
card_h = Inches(4.8)
gap = Inches(0.35)
start_x = Inches(0.6)
y = Inches(1.9)

for i, (name, price, accent, features) in enumerate(plans):
    x = start_x + (card_w + gap) * i
    b = box(slide, x, y, card_w, card_h, WHITE)

    # Popular badge for Club
    if name == "Club":
        badge_w = Inches(1.4)
        badge_h = Inches(0.28)
        badge = box(slide, x + (card_w - badge_w) / 2, y - Inches(0.14), badge_w, badge_h, AMBER, radius=0.3)
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = "⭐ Más vendido"
        run.font.size = Pt(8)
        run.font.color.rgb = WHITE
        run.font.bold = True

    accent_line(slide, x + Inches(0.3), y + Inches(0.15), Inches(3.0), accent)
    text(slide, x, y + Inches(0.3), card_w, Inches(0.4),
         name, size=20, color=accent, bold=True, align=PP_ALIGN.CENTER)
    text(slide, x, y + Inches(0.7), card_w, Inches(0.6),
         price, size=32, color=DARK, bold=True, align=PP_ALIGN.CENTER)

    for j, feat in enumerate(features):
        text(slide, x + Inches(0.4), y + Inches(1.5) + Inches(j * 0.45), Inches(3.0), Inches(0.35),
             f"✓  {feat}", size=11, color=MUTED)

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 6 — Funcionalidades principales
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, LIGHT_BG)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.5), "FUNCIONALIDADES")
text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
     "Módulos de la plataforma", size=34, color=DARK, bold=True)

modules = [
    ("👥", "Alumnos", "Alta, baja, edición.\nFicha completa con nivel,\ncontacto y pagos.", PRIMARY),
    ("📅", "Calendario", "Vista semana/mes.\nArrastrar sesiones.\nPanel lateral detallado.", EMERALD),
    ("📋", "Clases", "Horarios por día.\nGestión de inscripciones.\nControl de plazas.", VIOLET),
    ("🔄", "Recuperaciones", "Generación automática.\nControl de vencimientos.\nReubicación de alumnos.", AMBER),
    ("📊", "Dashboard", "KPIs en tiempo real.\nGráficos y tendencias.\n3 niveles de detalle.", RED),
    ("📱", "Notificaciones", "Email automático.\nIntegración WhatsApp.\nPlantillas personalizadas.", RGBColor(0x06, 0xB6, 0xD4)),
]

card_w = Inches(3.8)
card_h = Inches(2.6)
gap_x = Inches(0.35)
gap_y = Inches(0.4)
start_x = Inches(0.6)
start_y = Inches(1.9)

for i, (emoji, title, desc, accent) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = start_x + (card_w + gap_x) * col
    y = start_y + (card_h + gap_y) * row

    b = box(slide, x, y, card_w, card_h, WHITE)
    accent_line(slide, x + Inches(0.3), y + Inches(0.12), Inches(1.5), accent)
    text(slide, x + Inches(0.3), y + Inches(0.25), Inches(3), Inches(0.4),
         f"{emoji}  {title}", size=16, color=DARK, bold=True)
    text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.8),
         desc, size=11, color=MUTED)

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 7 — KPIs / Datos Reales
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.5), "DATOS REALES")
text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
     "El sistema en cifras", size=34, color=WHITE, bold=True)

# KPI row 1
kpis_top = [
    (str(DATA["totalAlumnos"]), "Alumnos activos", EMERALD),
    (str(DATA["totalClases"]), "Clases activas", PRIMARY),
    (str(DATA["profesores"]), "Profesores", VIOLET),
    (str(DATA["pistas"]), "Pistas", AMBER),
]

x_start = Inches(0.6)
for i, (val, label, accent) in enumerate(kpis_top):
    x = x_start + Inches(i * 3.1)
    y = Inches(2.0)
    card_w_kpi = Inches(2.8)
    card_h_kpi = Inches(1.6)
    b = box(slide, x, y, card_w_kpi, card_h_kpi, RGBColor(0x16, 0x1E, 0x30))
    accent_line(slide, x + Inches(0.3), y + Inches(0.12), Inches(2.0), accent)
    text(slide, x, y + Inches(0.3), card_w_kpi, Inches(0.6), val,
         size=36, color=accent, bold=True, align=PP_ALIGN.CENTER)
    text(slide, x, y + Inches(1.0), card_w_kpi, Inches(0.4), label,
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

# KPI row 2
kpis_bot = [
    (f"{DATA['mrr']}€", "MRR Mensual", EMERALD),
    (DATA["ocupacion"], "Ocupación", PRIMARY),
    (DATA["retencion"], "Retención", VIOLET),
    (str(DATA["recuperaciones"]), "Recuperaciones pendientes", AMBER),
]

for i, (val, label, accent) in enumerate(kpis_bot):
    x = x_start + Inches(i * 3.1)
    y = Inches(4.0)
    card_w_kpi = Inches(2.8)
    card_h_kpi = Inches(1.6)
    b = box(slide, x, y, card_w_kpi, card_h_kpi, RGBColor(0x16, 0x1E, 0x30))
    accent_line(slide, x + Inches(0.3), y + Inches(0.12), Inches(2.0), accent)
    text(slide, x, y + Inches(0.3), card_w_kpi, Inches(0.6), val,
         size=36, color=accent, bold=True, align=PP_ALIGN.CENTER)
    text(slide, x, y + Inches(1.0), card_w_kpi, Inches(0.4), label,
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4),
     "Datos extraídos del sistema en producción · Marzo 2026",
     size=10, color=RGBColor(0x4A, 0x5A, 0x6A), align=PP_ALIGN.CENTER)

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 8 — Cuadro de Clases
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, LIGHT_BG)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.5), "CLASES")
text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
     "Horario semanal de clases", size=34, color=DARK, bold=True)

# Table header
headers = ["Clase", "Horario", "Pista", "Profesor", "Ocupación"]
col_widths = [Inches(2.8), Inches(3.0), Inches(1.5), Inches(1.8), Inches(1.5)]
x_start = Inches(0.8)
y_header = Inches(2.0)

# Header row bg
rect(slide, x_start, y_header, sum(cw for cw in col_widths) + Inches(0.4), Inches(0.45), PRIMARY)
x = x_start + Inches(0.2)
for i, h in enumerate(headers):
    text(slide, x, y_header + Inches(0.05), col_widths[i], Inches(0.35),
         h, size=11, color=WHITE, bold=True)
    x += col_widths[i]

# Data rows
row_colors = [WHITE, SURFACE]
for r, (nombre, horario, pista, prof, ocup) in enumerate(DATA["clases"]):
    y = y_header + Inches(0.5) + Inches(r * 0.5)
    row_bg = row_colors[r % 2]
    rect(slide, x_start, y, sum(cw for cw in col_widths) + Inches(0.4), Inches(0.45), row_bg)

    vals = [nombre, horario, pista, prof, ocup]
    x = x_start + Inches(0.2)
    for i, v in enumerate(vals):
        c = DARK
        b = False
        if i == 4:  # Occupancy
            nums = v.split("/")
            if int(nums[0]) >= int(nums[1]) - 1:
                c = RED
            elif int(nums[0]) == 0:
                c = MUTED
            else:
                c = EMERALD
            b = True
        text(slide, x, y + Inches(0.05), col_widths[i], Inches(0.35),
             v, size=10, color=c, bold=b)
        x += col_widths[i]

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 9 — Screenshot: Dashboard Elite
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.4), "CAPTURA")
text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
     "Dashboard Elite", size=28, color=WHITE, bold=True)

# Screenshot
add_screenshot(slide, "dashboard_elite.png", Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 10 — Screenshot: Calendario
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.4), "CAPTURA")
text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
     "Calendario", size=28, color=WHITE, bold=True)

add_screenshot(slide, "calendario.png", Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 11 — Screenshot: Clases
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.4), "CAPTURA")
text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
     "Gestión de Clases", size=28, color=WHITE, bold=True)

add_screenshot(slide, "clases.png", Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 12 — Screenshot: Alumnos
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.4), "CAPTURA")
text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
     "Gestión de Alumnos", size=28, color=WHITE, bold=True)

add_screenshot(slide, "alumnos.png", Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 13 — Screenshot: Planes
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.4), "CAPTURA")
text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
     "Selección de Planes", size=28, color=WHITE, bold=True)

add_screenshot(slide, "planes.png", Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 14 — Screenshots: Recuperaciones + Notificaciones
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.4), "CAPTURAS")
text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
     "Recuperaciones & Notificaciones", size=28, color=WHITE, bold=True)

add_screenshot(slide, "recuperaciones.png", Inches(0.3), Inches(1.6), Inches(6.2), Inches(5.4))
add_screenshot(slide, "notificaciones.png", Inches(6.8), Inches(1.6), Inches(6.2), Inches(5.4))

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 15 — Comparativa de Dashboards
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.4), "CAPTURAS")
text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
     "Dashboard por Plan", size=28, color=WHITE, bold=True)

# 3 screenshots side by side
scr_w = Inches(4.0)
scr_h = Inches(4.8)
gap = Inches(0.3)

labels = [("Starter", EMERALD), ("Club", AMBER), ("Elite", RED)]
files = ["dashboard_starter.png", "dashboard_club.png", "dashboard_elite.png"]
x = Inches(0.4)
for i, (label, color) in enumerate(labels):
    add_screenshot(slide, files[i], x, Inches(2.0), scr_w, scr_h)
    text(slide, x, Inches(1.55), scr_w, Inches(0.4),
         label, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    x += scr_w + gap

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 16 — Roadmap / Próximos pasos
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, LIGHT_BG)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)

section_tag(slide, Inches(0.8), Inches(0.5), "ROADMAP")
text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
     "Próximos pasos", size=34, color=DARK, bold=True)

phases = [
    ("Q2 2026", "Mejoras UI/UX", [
        "Vista diaria del club",
        "Drag & drop en calendario",
        "App móvil responsive",
        "Tema oscuro",
    ], EMERALD),
    ("Q3 2026", "Integraciones", [
        "Pasarela de pagos (Stripe)",
        "WhatsApp Business API",
        "Google Calendar sync",
        "Exportación a Excel/PDF",
    ], PRIMARY),
    ("Q4 2026", "Escalabilidad", [
        "Multi-club con datos separados",
        "API pública documentada",
        "Portal para familias",
        "Onboarding automatizado",
    ], VIOLET),
]

card_w = Inches(3.8)
card_h = Inches(3.8)
gap = Inches(0.35)
x_start = Inches(0.6)

for i, (quarter, title, items_list, accent) in enumerate(phases):
    x = x_start + (card_w + gap) * i
    y = Inches(2.0)
    b = box(slide, x, y, card_w, card_h, WHITE)
    accent_line(slide, x + Inches(0.3), y + Inches(0.12), Inches(1.5), accent)
    text(slide, x + Inches(0.3), y + Inches(0.25), Inches(1.5), Inches(0.3),
         quarter, size=11, color=accent, bold=True)
    text(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.2), Inches(0.4),
         title, size=18, color=DARK, bold=True)
    bullets(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.2), Inches(2.5),
            items_list, size=11, color=MUTED)

add_footer(slide)


# ═══════════════════════════════════════════════
# SLIDE 17 — Cierre
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_bg(slide, DARK)
rect(slide, 0, 0, W, Inches(0.05), PRIMARY)
rect(slide, 0, Inches(0.05), W, Inches(0.03), AMBER)

text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
     "Academia Pádel", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
     "Gestión inteligente para tu club de pádel",
     size=22, color=PRIMARY, align=PP_ALIGN.CENTER)

accent_line(slide, Inches(4.5), Inches(4.1), Inches(4.3), AMBER)

text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
     "¿Preguntas?",
     size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
     "Grupo ALTER  ·  Departamento de Informática  ·  Marzo 2026",
     size=14, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(slide)


# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
prs.save(OUTPUT)
print(f"✅ Presentación guardada en: {OUTPUT}")
print(f"   {len(prs.slides)} diapositivas generadas")
