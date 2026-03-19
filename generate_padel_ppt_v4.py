#!/usr/bin/env python3
"""
Presentación COMERCIAL «Academia Pádel» — Plantilla Padel Tennis.
Enfoque: Venta a dueños de club que usan Excel.
20 diapositivas con el estilo de la plantilla deportiva.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os

TEMPLATE = "/docker/data_dockers/Padel Tennis Powerpoint.pptx"
SS       = "/docker/data_dockers/padel_app/screenshots"
OUTPUT   = "/docker/data_dockers/padel_app/Academia_Padel_Presentacion.pptx"

prs = Presentation(TEMPLATE)

# ── helpers ──────────────────────────────────

def _set(shape, txt):
    """Replace ALL text in shape (clears every run, writes to the first one)."""
    if not shape.has_text_frame:
        return
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    for r in runs:
        r.text = ""
    if runs:
        runs[0].text = txt
    else:
        p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
        p.add_run().text = txt

def _pic(slide, idx, path):
    """Insert image into picture placeholder with center-crop."""
    if not os.path.exists(path):
        print(f"  ⚠ {path}"); return
    for sh in slide.placeholders:
        if sh.placeholder_format.idx == idx:
            pw, ph_ = sh.width, sh.height
            img = Image.open(path)
            iw, ih = img.size
            pic = sh.insert_picture(path)
            pr = pw / ph_; ir = iw / ih
            if abs(ir - pr) > 0.05:
                if ir > pr:
                    c = (iw - ih * pr) / (2 * iw)
                    pic.crop_left = c; pic.crop_right = c
                else:
                    c = (ih - iw / pr) / (2 * ih)
                    pic.crop_top = c; pic.crop_bottom = c
            return
    print(f"  ⚠ PH idx={idx}")

def _tx(slide, fragment):
    """Find shapes containing text fragment."""
    return [s for s in slide.shapes if hasattr(s, "text") and fragment in s.text]

# ══════════════════════════════════════════════
# SLIDE 1 — PORTADA
# Layout: 1_Title Slide (full bg PH idx=10)
# "FAST / SOCIAL / GROWING" + tagline
# ══════════════════════════════════════════════
print("Sl 1: Portada")
s = prs.slides[0]
_pic(s, 10, f"{SS}/dashboard_elite_full.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if t == "FAST":           _set(sh, "FÁCIL.")
    elif t == "SOCIAL.":      _set(sh, "COMPLETA.")
    elif t == "GROWING":      _set(sh, "RENTABLE.")
    elif "BUILDING THE" in t: _set(sh, "DEJA EL EXCEL. DIGITALIZA TU CLUB")
    elif "PADEL TENNIS PRESENTATION" in t: _set(sh, "ACADEMIA PÁDEL — GESTIÓN INTELIGENTE")

# ══════════════════════════════════════════════
# SLIDE 2 — EL PROBLEMA
# Layout: 2_Title Slide (image PH idx=10 bottom-right, stat card)
# ══════════════════════════════════════════════
print("Sl 2: El Problema")
s = prs.slides[1]
_pic(s, 10, f"{SS}/dashboard_elite.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Where Sport" in t:        _set(sh, "¿Tu academia aún funciona con Excel?")
    elif "Padel tennis is not" in t: _set(sh, "Alumnos en una hoja, clases en otra, cobros en una libreta… Si alguien toca la fórmula equivocada, pierdes los datos del mes.")
    elif "65,5%" in t:            _set(sh, "73%")
    elif "Building the future" in t: _set(sh, "de clubes tienen errores cada mes en sus hojas de cálculo")

# ══════════════════════════════════════════════
# SLIDE 3 — LA SOLUCIÓN
# Layout: 3_Title Slide (2 image PHs: idx=10 right, idx=11 bottom)
# ══════════════════════════════════════════════
print("Sl 3: La Solución")
s = prs.slides[2]
_pic(s, 10, f"{SS}/planes.png")
_pic(s, 11, f"{SS}/dashboard_club.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "The Sport Everyone" in t: _set(sh, "Todo tu club en una sola pantalla")
    elif "Padel tennis is not" in t: _set(sh, "Alumnos, clases, pagos, calendario y comunicación. Sin complicaciones, sin fórmulas, sin errores.")
    elif "BUILDING THE" in t:     _set(sh, "SI SABES USAR WHATSAPP, SABES USAR ACADEMIA PÁDEL")

# ══════════════════════════════════════════════
# SLIDE 4 — POR QUÉ AHORA
# Layout: 4_Title Slide (image PH idx=10 right + stat)
# ══════════════════════════════════════════════
print("Sl 4: Por qué ahora")
s = prs.slides[3]
_pic(s, 10, f"{SS}/alumnos.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "The Global Boom" in t:    _set(sh, "El pádel no para de crecer. Tu gestión también debe hacerlo")
    elif t == "65,5%":            _set(sh, "+5h")
    elif "Building the future" in t: _set(sh, "horas libres a la semana para dedicar a tus alumnos")
    elif "Padel tennis is not" in t: _set(sh, "Más de 8 millones de jugadores en España. Las academias que digitalizan su gestión retienen un 30% más de alumnos.")
    elif "Families, athletes" in t:
        if sh.top < Inches(5.5):
            _set(sh, "Desde el móvil, tablet o PC: acceso en cualquier momento.")
        else:
            _set(sh, "Sin instalación, sin cursos. Funcionando en 5 minutos.")

# ══════════════════════════════════════════════
# SLIDE 5 — CALENDARIO
# Layout: 5_Title Slide (image PH idx=10 right)
# ══════════════════════════════════════════════
print("Sl 5: Calendario")
s = prs.slides[4]
_pic(s, 10, f"{SS}/calendario.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "The Sport Everyone" in t: _set(sh, "Calendario semanal: todo de un vistazo")
    elif "Padel tennis is not" in t: _set(sh, "Ves todas las clases de la semana: quién da clase, en qué pista y a qué hora. El sistema evita solapamientos automáticamente.")
    elif "How we stand out" in t: _set(sh, "Nunca más dobles reservas ni confusiones de horario.")
    elif "BUILDING THE" in t:     _set(sh, "NAVEGA ENTRE SEMANAS CON UN SOLO CLIC")

# ══════════════════════════════════════════════
# SLIDE 6 — DATOS DEL CLUB (con gráfico)
# Layout: 8_Title Slide (image PH idx=10 left + stat group)
# ══════════════════════════════════════════════
print("Sl 6: Datos del Club")
s = prs.slides[5]
_pic(s, 10, f"{SS}/dashboard_elite.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "65,5%" in t:              _set(sh, "826€")
    elif "Building the future" in t: _set(sh, "MRR mensual con solo 14 alumnos activos")

# ══════════════════════════════════════════════
# SLIDE 7 — EXCEL vs APP (gráfico de barras)
# Layout: Title Slide (chart + text)
# ══════════════════════════════════════════════
print("Sl 7: Excel vs App")
s = prs.slides[6]
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Making Noise" in t:       _set(sh, "Excel vs Academia Pádel")
    elif "Padel tennis is not" in t: _set(sh, "El Excel no avisa de impagos, no gestiona faltas, no envía notificaciones. Academia Pádel lo hace todo automáticamente.")
    elif "Families, athletes" in t: _set(sh, "Avisos de impago, control de asistencia, comunicación directa.")
    elif "BUILDING THE" in t:     _set(sh, "LA COMPARACIÓN QUE LO CAMBIA TODO")

# ══════════════════════════════════════════════
# SLIDE 8 — CLASES (2 imágenes)
# Layout: 6_Title Slide (2 image PHs: idx=10 left, idx=11 center)
# ══════════════════════════════════════════════
print("Sl 8: Clases")
s = prs.slides[7]
_pic(s, 10, f"{SS}/clases.png")
_pic(s, 11, f"{SS}/calendario.png")
fdone = 0
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Families, athletes" in t:
        fdone += 1
        if fdone == 1:
            _set(sh, "Crea grupos, asigna profesor y pista en 1 clic.")
        else:
            _set(sh, "Control de ocupación y asistencia automática.")

# ══════════════════════════════════════════════
# SLIDE 9 — PERFIL / TESTIMONIAL
# Layout: 7_Title Slide (3 image PHs + profile info)
# ══════════════════════════════════════════════
print("Sl 9: Testimonial")
s = prs.slides[8]
_pic(s, 10, f"{SS}/alumnos.png")
_pic(s, 11, f"{SS}/recuperaciones.png")
_pic(s, 12, f"{SS}/notificaciones.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Rosa Maria" in t:         _set(sh, "Director de Club")
    elif "Padel tennis is not" in t: _set(sh, '"Desde que dejamos el Excel, no he vuelto a perder un cobro ni a olvidar una recuperación. Ojalá lo hubiera tenido antes."')
    elif t == "Skills":           _set(sh, "Facilidad")
    elif t == "Experience":       _set(sh, "Ahorro tiempo")
    elif t == "90%":              _set(sh, "95%")
    elif t == "80%":              _set(sh, "90%")
    elif t == "2026":             _set(sh, "★★★★★")

# ══════════════════════════════════════════════
# SLIDE 10 — MÉTRICAS EN ACCIÓN (3 images + stats)
# Layout: 9_Title Slide (3 image PHs)
# ══════════════════════════════════════════════
print("Sl 10: Métricas")
s = prs.slides[9]
_pic(s, 10, f"{SS}/dashboard_elite.png")
_pic(s, 11, f"{SS}/dashboard_club.png")
_pic(s, 12, f"{SS}/dashboard_starter.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Making Noise" in t:       _set(sh, "Métricas reales de la demo")
    elif "Padel tennis is not" in t: _set(sh, "14 alumnos · 8 clases · 3 profesores · 3 pistas · Ocupación 44% · Retención 94%")
    elif "Families, athletes" in t: _set(sh, "Datos que se actualizan automáticamente.")
    elif "$672,9" in t:           _set(sh, "826€ MRR")

# ══════════════════════════════════════════════
# SLIDE 11 — ALUMNOS (image right)
# Layout: 10_Title Slide (image PH idx=10)
# ══════════════════════════════════════════════
print("Sl 11: Alumnos")
s = prs.slides[10]
_pic(s, 10, f"{SS}/alumnos.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Collaborating" in t:      _set(sh, "Ficha completa de cada alumno")
    elif "Padel tennis is not" in t: _set(sh, "Datos de contacto, nivel, grupo, asistencia y estado de pago en una sola pantalla. Sabes al instante quién ha pagado y quién no.")

# ══════════════════════════════════════════════
# SLIDE 12 — SCREENSHOT FULL (decorativo)
# Layout: 11_Title Slide (image PH idx=10)
# ══════════════════════════════════════════════
print("Sl 12: Vista Completa")
s = prs.slides[11]
_pic(s, 10, f"{SS}/recuperaciones.png")

# ══════════════════════════════════════════════
# SLIDE 13 — RECUPERACIONES
# Layout: 12_Title Slide (image PH idx=10 right)
# ══════════════════════════════════════════════
print("Sl 13: Recuperaciones")
s = prs.slides[12]
_pic(s, 10, f"{SS}/recuperaciones.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Play safe" in t:          _set(sh, "Faltas y recuperaciones automáticas")
    elif "Padel tennis is not" in t: _set(sh, "Cuando un alumno falta, el sistema registra la ausencia y le ofrece recuperar en otro horario. Sin llamadas, sin post-its, sin olvidos.")
    elif "BUILDING THE" in t:     _set(sh, "0 MINUTOS DE GESTIÓN POR RECUPERACIÓN")

# ══════════════════════════════════════════════
# SLIDE 14 — NOTIFICACIONES (image top-right + stats)
# Layout: 13_Title Slide (image PH idx=10)
# ══════════════════════════════════════════════
print("Sl 14: Notificaciones")
s = prs.slides[13]
_pic(s, 10, f"{SS}/notificaciones.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Collaborating" in t:      _set(sh, "Comunicación directa con tus alumnos")
    elif "Padel tennis is not" in t: _set(sh, "Envía avisos de clase, cambios de horario o recordatorios de pago desde la app. El alumno recibe la notificación al instante.")
    elif t == "65,5%":            _set(sh, "100%")
    elif "Building the future" in t: _set(sh, "de notificaciones entregadas al instante")
    elif "$672,9" in t:           _set(sh, "0 llamadas")

# ══════════════════════════════════════════════
# SLIDE 15 — GRÁFICO DE EVOLUCIÓN
# Layout: Title Slide (chart + stats)
# ══════════════════════════════════════════════
print("Sl 15: Evolución")
s = prs.slides[14]
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Play safe" in t:          _set(sh, "Crecimiento mes a mes")
    elif "Padel tennis is not" in t: _set(sh, "Los clubes que digitalizan su gestión ven resultados desde el primer mes: menos errores, más retención, mejor control económico.")
    elif "$672,9" in t:           _set(sh, "+30%")
    elif "65,5%" in t:            _set(sh, "94%")
    elif "Building the future" in t: _set(sh, "Tasa de retención de alumnos")

# ══════════════════════════════════════════════
# SLIDE 16 — DOS SCREENSHOTS (2 image PHs)
# Layout: 14_Title Slide (2 images right side)
# ══════════════════════════════════════════════
print("Sl 16: Dos Pantallas")
s = prs.slides[15]
_pic(s, 10, f"{SS}/clases.png")
_pic(s, 11, f"{SS}/calendario.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Making Noise" in t:       _set(sh, "Clases y Calendario conectados")
    elif "Padel tennis is not" in t: _set(sh, "Creas una clase y aparece automáticamente en el calendario. Modificas el horario y se actualiza en todas las vistas.")
    elif "BUILDING THE" in t:     _set(sh, "TODO SINCRONIZADO EN TIEMPO REAL")

# ══════════════════════════════════════════════
# SLIDE 17 — PLANES
# Layout: 15_Title Slide (2 image PHs)
# ══════════════════════════════════════════════
print("Sl 17: Planes")
s = prs.slides[16]
_pic(s, 10, f"{SS}/planes.png")
_pic(s, 11, f"{SS}/dashboard_starter.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Collaborating" in t:      _set(sh, "Un plan para cada academia")
    elif "Padel tennis is not" in t: _set(sh, "Empieza con lo básico y crece cuando quieras. Sin permanencia, sin letra pequeña. Starter 29€ · Club 49€ · Elite 79€")
    elif "$672,9" in t:           _set(sh, "Desde 29€")

# ══════════════════════════════════════════════
# SLIDE 18 — GALERÍA (3 images)
# Layout: 17_Title Slide (3 image PHs)
# ══════════════════════════════════════════════
print("Sl 18: Galería")
s = prs.slides[17]
_pic(s, 10, f"{SS}/dashboard_elite.png")
_pic(s, 11, f"{SS}/alumnos.png")
_pic(s, 12, f"{SS}/clases.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Making Noise" in t:       _set(sh, "Así se ve por dentro")
    elif "Padel tennis is not" in t: _set(sh, "Dashboard con KPIs reales, fichas de alumnos con control de pagos, y gestión de clases con ocupación en tiempo real.")
    elif "90%" in t:              _set(sh, "14 días")

# ══════════════════════════════════════════════
# SLIDE 19 — CONTACTO / CTA
# Layout: 16_Title Slide (wide image PH top + contact)
# ══════════════════════════════════════════════
print("Sl 19: Contacto")
s = prs.slides[18]
_pic(s, 10, f"{SS}/planes.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Connect" in t and "Us" in t: _set(sh, "¿Lo probamos juntos?")
    elif "Padel tennis is not" in t: _set(sh, "Agenda una demo de 15 min con los datos de tu club.")
    elif "hello@yourmail" in t:   _set(sh, "info@academiapadel.com")
    elif "+123-456-7890" in t:    _set(sh, "+34 912 345 678")
    elif "www.yourwebsite" in t:  _set(sh, "app.academiapadel.com")
    elif "BUILDING THE" in t:     _set(sh, "14 DÍAS GRATIS · SIN TARJETA · SIN COMPROMISO")

# ══════════════════════════════════════════════
# SLIDE 20 — CIERRE
# Layout: 1_Title Slide (full bg + thank you)
# ══════════════════════════════════════════════
print("Sl 20: Cierre")
s = prs.slides[19]
_pic(s, 10, f"{SS}/dashboard_elite_full.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "THANK YOU" in t:          _set(sh, "GRACIAS.")
    elif "FOR JOINING" in t:      _set(sh, "TU ACADEMIA MERECE MÁS QUE UN EXCEL")
    elif "BUILDING THE" in t:     _set(sh, "APP.ACADEMIAPADEL.COM")

# ──────────────────────────────────────────────
prs.save(OUTPUT)
sz = os.path.getsize(OUTPUT) / 1024 / 1024
print(f"\n✅ {OUTPUT}")
print(f"   {len(prs.slides)} diapositivas | {sz:.1f} MB")
