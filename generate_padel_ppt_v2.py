#!/usr/bin/env python3
"""
Genera la presentación PPT «Academia Pádel — SaaS de Gestión»
reutilizando las diapositivas de la plantilla 'Ai Automation Business ppt.pptx'
y reemplazando textos + insertando screenshots en los Picture Placeholders.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import copy, os

# ──────────────────────────────────────────────
# Paths
# ──────────────────────────────────────────────
TEMPLATE   = "/docker/data_dockers/Ai Automation Business ppt.pptx"
SCREENSHOTS = "/docker/data_dockers/padel_app/screenshots"
OUTPUT     = "/docker/data_dockers/padel_app/Academia_Padel_Presentacion.pptx"

# ──────────────────────────────────────────────
# Load template
# ──────────────────────────────────────────────
prs = Presentation(TEMPLATE)
W = prs.slide_width    # 13.33"
H = prs.slide_height   # 7.50"


# ──────────────────────────────────────────────
# Helper: find a shape by name pattern in a slide
# ──────────────────────────────────────────────
def find_shape(slide, name_fragment):
    for s in slide.shapes:
        if name_fragment in s.name:
            return s
    return None

def find_shapes(slide, name_fragment):
    return [s for s in slide.shapes if name_fragment in s.name]

def set_text_in_shape(shape, new_text, font_size=None, bold=None, color=None):
    """Replace all text in a shape preserving formatting where possible."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Use first paragraph's first run as the style reference
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # Set text on the first run of first paragraph
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
        if font_size is not None:
            tf.paragraphs[0].runs[0].font.size = Pt(font_size)
        if bold is not None:
            tf.paragraphs[0].runs[0].font.bold = bold
        if color is not None:
            tf.paragraphs[0].runs[0].font.color.rgb = color
    else:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = new_text
        if font_size:
            run.font.size = Pt(font_size)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color

def set_multiline(shape, lines, font_size=None, bold=None, color=None):
    """Set lines preserving first run style, using line breaks."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear existing runs
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # Find the first run to use as template
    first_run = None
    for para in tf.paragraphs:
        for run in para.runs:
            first_run = run
            break
        if first_run:
            break
    if first_run:
        first_run.text = "\n".join(lines)
        if font_size:
            first_run.font.size = Pt(font_size)
        if bold is not None:
            first_run.font.bold = bold
        if color:
            first_run.font.color.rgb = color

def set_text_tworuns(shape, line1, line2, size1=None, size2=None):
    """For shapes that had two runs (like title + subtitle on same shape)."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    runs = []
    for para in tf.paragraphs:
        for run in para.runs:
            runs.append(run)
    # Clear all
    for r in runs:
        r.text = ""
    if len(runs) >= 2:
        runs[0].text = line1
        runs[1].text = line2
        if size1:
            runs[0].font.size = Pt(size1)
        if size2:
            runs[1].font.size = Pt(size2)
    elif len(runs) == 1:
        runs[0].text = f"{line1}\n{line2}"

def insert_screenshot(slide, placeholder_idx, image_path):
    """Insert a screenshot into a picture placeholder, preserving aspect ratio."""
    if not os.path.exists(image_path):
        print(f"  ⚠ Image not found: {image_path}")
        return
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            shape.insert_picture(image_path)
            return
    print(f"  ⚠ Placeholder idx={placeholder_idx} not found")

def crop_screenshot_to_fit(placeholder_w, placeholder_h, img_w, img_h):
    """Calculate crop values to fit image in placeholder while maintaining center.
    Returns (crop_left, crop_top, crop_right, crop_bottom) as fractions."""
    ph_ratio = placeholder_w / placeholder_h
    img_ratio = img_w / img_h
    
    if img_ratio > ph_ratio:
        # Image is wider than placeholder → crop sides
        visible_w = img_h * ph_ratio
        crop_x = (img_w - visible_w) / (2 * img_w)
        return (crop_x, 0, crop_x, 0)
    else:
        # Image is taller than placeholder → crop top/bottom
        visible_h = img_w / ph_ratio
        crop_y = (img_h - visible_h) / (2 * img_h)
        return (0, crop_y, 0, crop_y)

def insert_screenshot_cropped(slide, placeholder_idx, image_path):
    """Insert screenshot with proper cropping to fill placeholder without distortion."""
    if not os.path.exists(image_path):
        print(f"  ⚠ Image not found: {image_path}")
        return
    
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            # Get placeholder dimensions 
            ph_w = shape.width
            ph_h = shape.height
            
            # Get image dimensions
            img = Image.open(image_path)
            img_w, img_h = img.size
            
            # Insert picture
            pic = shape.insert_picture(image_path)
            
            # Calculate and apply cropping
            ph_ratio = ph_w / ph_h
            img_ratio = img_w / img_h
            
            if abs(img_ratio - ph_ratio) > 0.05:
                cl, ct, cr, cb = crop_screenshot_to_fit(ph_w, ph_h, img_w, img_h)
                pic.crop_left = cl
                pic.crop_top = ct
                pic.crop_right = cr
                pic.crop_bottom = cb
            
            return
    print(f"  ⚠ Placeholder idx={placeholder_idx} not found")


# ══════════════════════════════════════════════
# SLIDE 1 — PORTADA (layout: '1_Title Slide')
# Full-page Picture PH idx=10, title text, subtitle
# ══════════════════════════════════════════════
print("Slide 1: Portada")
s = prs.slides[0]
# Insert dashboard screenshot as bg image
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/dashboard_elite_full.png")
# Replace "Automation\nSTRATEGY" → "Academia\nPÁDEL"
for shape in s.shapes:
    if hasattr(shape, "text"):
        if "Automation" in shape.text:
            set_text_tworuns(shape, "Academia", "PÁDEL")
        elif "Embracing AI" in shape.text:
            set_text_in_shape(shape, "Gestión integral de tu academia de pádel · SaaS multi-plan")
        elif "Read More" in shape.text:
            set_text_in_shape(shape, "Descubrir")


# ══════════════════════════════════════════════
# SLIDE 2 — DASHBOARD ELITE (layout: 'Custom Layout')
# Image PH idx=10 (6.75"x3.97"), text boxes
# ══════════════════════════════════════════════
print("Slide 2: Dashboard Elite")
s = prs.slides[1]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/dashboard_elite.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        if "The Basics of" in shape.text:
            set_text_in_shape(shape, "Dashboard Elite — Vista General")
        elif "AI Business Strategy" == shape.text.strip():
            set_text_in_shape(shape, "Panel de Control")
        elif "AI Business Strategy leverages" in shape.text:
            set_text_in_shape(shape, "Dashboard adaptativo con KPIs en tiempo real: MRR, ocupación de pistas, retención de alumnos, gráficos de evolución mensual y distribución por nivel. La vista cambia según el plan contratado (Starter, Club, Elite).")
        elif "Read More" in shape.text:
            set_text_in_shape(shape, "Ver Demo")


# ══════════════════════════════════════════════
# SLIDE 3 — KPIs (layout: 'Title Slide')
# Stats layout with Q2/Q3/Q4 sections and metrics
# ══════════════════════════════════════════════
print("Slide 3: KPIs")
s = prs.slides[2]
# Map the stats: 201% → 14, 20M+ → 826€, 290+ → 94%
# Map quarters: Q2 → Alumnos, Q3 → MRR, Q4 → Retención
stat_map = {
    "201": "14",
    "20M": "826€",
    "290": "94%",
    "Predictive Business Models": "Alumnos Activos",
    "Workflow Automation Tools": "MRR Mensual",
    "Business Intelligence Tools": "Tasa Retención",
    "Q2": "📊",
    "Q3": "💰",
    "Q4": "🏆",
}
desc_replaced = False
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        for old, new in stat_map.items():
            if old in txt:
                set_text_in_shape(shape, new)
                break
        if "AI Business Strategy leverages" in txt and not desc_replaced:
            set_text_in_shape(shape, "Métricas capturadas en tiempo real desde la API de la plataforma. El sistema monitoriza ocupación de pistas (44%), distribución por nivel y evolución semanal de alumnos.")
            desc_replaced = True


# ══════════════════════════════════════════════
# SLIDE 4 — AGENDA (layout: '4_Title Slide')
# Left text + Right image PH idx=10
# ══════════════════════════════════════════════
print("Slide 4: Agenda")
s = prs.slides[3]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/planes.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "What" in txt and "Explore" in txt:
            set_text_tworuns(shape, "Contenido de", "la Presentación")
        elif txt == "Agenda":
            set_text_in_shape(shape, "Agenda")
        elif "Real-Time" in txt:
            set_text_in_shape(shape, "1. Visión General")
        elif "Scalable" in txt:
            set_text_in_shape(shape, "4. Módulos Funcionales")
        elif "AI Business Strategy leverages artificial intelligence and automation." == txt:
            # Two short descs next to the agenda items
            if shape.top < Inches(4):
                set_text_in_shape(shape, "2. Arquitectura Técnica\n3. Modelo SaaS (3 planes)")
            else:
                set_text_in_shape(shape, "5. KPIs y Datos Reales\n6. Roadmap 2026")
        elif "AI Business Strategy leverages" in txt:
            set_text_in_shape(shape, "Plataforma SaaS completa para la gestión integral de academias de pádel. Incluye gestión de alumnos, clases, calendario, recuperaciones, notificaciones y dashboards adaptados por plan.")


# ══════════════════════════════════════════════
# SLIDE 5 — CALENDARIO (layout: '2_Title Slide')
# Big image PH idx=11 (left), text sidebar right
# ══════════════════════════════════════════════
print("Slide 5: Calendario")
s = prs.slides[4]
insert_screenshot_cropped(s, 11, f"{SCREENSHOTS}/calendario.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "AI-Enhanced" in txt:
            set_text_in_shape(shape, "Vista Semanal Interactiva")
        elif "Customer Insight" in txt:
            set_text_in_shape(shape, "Navegación Temporal")
        elif "Operational Efficiency" in txt:
            set_text_in_shape(shape, "Calendario de Clases")
        # Update description texts
        for sub_shape in s.shapes:
            if hasattr(sub_shape, "text") and "AI Business Strategy leverages artificial intelligence and automation." == sub_shape.text.strip():
                # There are two of these, first = calendar desc, second = nav desc
                if "AI-Enhanced" in txt or "Customer" in txt:
                    continue
                break


# Fix the two description boxes on slide 5
desc_count = 0
for shape in s.shapes:
    if hasattr(shape, "text") and shape.text.strip() == "AI Business Strategy leverages artificial intelligence and automation.":
        desc_count += 1
        if desc_count == 1:
            set_text_in_shape(shape, "Visualización de todas las clases por semana con código de colores por nivel.")
        elif desc_count == 2:
            set_text_in_shape(shape, "Selector de semanas anterior/siguiente con vista rápida del horario.")


# ══════════════════════════════════════════════
# SLIDE 6 — CLASES (layout: '5_Title Slide')
# Image PH idx=10 (right side), text left
# ══════════════════════════════════════════════
print("Slide 6: Clases")
s = prs.slides[5]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/clases.png")
stat_map6 = {
    "290+": "8",
    "250 MW": "3",
    "Automated Insights": "Clases Activas",
    "Predictive Analytics": "Profesores",
}
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "Identifying" in txt and "Use Cases" in txt:
            set_text_tworuns(shape, "Gestión de", "Clases")
        elif "Where to Apply" in txt:
            set_text_in_shape(shape, "Módulo de Clases")
        elif "AI Business Strategy leverages" in txt:
            set_text_in_shape(shape, "Listado completo de clases con filtrado por nivel, profesor y día. Cada clase muestra ocupación, horario y permite gestión de alumnos inscritos.")
        for old, new in stat_map6.items():
            if txt == old:
                set_text_in_shape(shape, new)
                break


# ══════════════════════════════════════════════
# SLIDE 7 — ALUMNOS (layout: '6_Title Slide')
# Wide image PH idx=10 at bottom
# ══════════════════════════════════════════════
print("Slide 7: Alumnos")
s = prs.slides[6]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/alumnos.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "Strategic Intelligence" == txt:
            set_text_in_shape(shape, "Gestión de Alumnos")
        elif "Innovation Model" == txt:
            set_text_in_shape(shape, "Fichas Individuales")
        elif "AI Business Strategy leverages artificial." == txt:
            if shape.top < Inches(2.5):
                set_text_in_shape(shape, "14 alumnos, búsqueda y filtrado por nivel.")
            else:
                set_text_in_shape(shape, "Alta, baja, edición, historial de pagos.")


# ══════════════════════════════════════════════
# SLIDE 8 — FULL SCREENSHOT (layout: '11_Title Slide')
# Big image PH idx=10 + decorative groups
# ══════════════════════════════════════════════
print("Slide 8: Dashboard Full")
s = prs.slides[7]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/planes.png")


# ══════════════════════════════════════════════
# SLIDE 9 — ARQUITECTURA (layout: '7_Title Slide')
# Wide image PH idx=10 at top, text below
# ══════════════════════════════════════════════
print("Slide 9: Arquitectura")
s = prs.slides[8]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/dashboard_elite_full.png")
stat_map9 = {
    "789+": "v15",
    "250 GBZ": "v5+",
}
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "Machine" in txt and "Learning" in txt:
            set_text_tworuns(shape, "Arquitectura", "Técnica")
        elif "Teaching Machines" in txt:
            set_text_in_shape(shape, "Stack Tecnológico")
        elif "AI Business Strategy leverages artificial intelligence and automation technologies" in txt:
            set_text_in_shape(shape, "Frontend: Next.js 15 + TypeScript + Tailwind v4\nBackend: Express v5 + Prisma v6\nBase de datos: PostgreSQL 16\nInfra: Docker Compose multi-contenedor")
        for old, new in stat_map9.items():
            if txt == old:
                set_text_in_shape(shape, new)
                break
        if "AI Business Strategy leverages artificial intelligence." == txt:
            if shape.left < Inches(7):
                set_text_in_shape(shape, "Next.js 15, React 19, TypeScript, Tailwind CSS")
            else:
                set_text_in_shape(shape, "Express v5, Prisma ORM, PostgreSQL 16")


# ══════════════════════════════════════════════
# SLIDE 10 — COMPARATIVA 3 DASHBOARDS (layout: '8_Title Slide')
# 3 image PHs: idx=10, 11, 12 side by side
# ══════════════════════════════════════════════
print("Slide 10: Comparativa Dashboards")
s = prs.slides[9]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/dashboard_starter.png")
insert_screenshot_cropped(s, 11, f"{SCREENSHOTS}/dashboard_club.png")
insert_screenshot_cropped(s, 12, f"{SCREENSHOTS}/dashboard_elite.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "Developing" in txt and "Strategy" in txt:
            set_text_tworuns(shape, "Comparativa", "Dashboards")
        elif "Aligning AI" in txt:
            set_text_in_shape(shape, "Vista por Plan")
        elif "AI Business Strategy leverages artificial intelligence and automation technologies to optimize." == txt:
            set_text_in_shape(shape, "Cada plan muestra un dashboard diferenciado con métricas y gráficos adaptados al nivel de suscripción.")
        elif "Read More" in txt:
            set_text_in_shape(shape, "Ver Demo")
        elif "Machine Learning" == txt:
            set_text_in_shape(shape, "Plan Starter")
        elif "Business Automation" == txt:
            set_text_in_shape(shape, "Plan Club")
        elif "Intelligent Platform" == txt:
            set_text_in_shape(shape, "Plan Elite")
        # The "AI Business Strategy" labels under each image
        elif txt == "AI Business Strategy":
            # Determine by position which one it is
            if shape.left < Inches(6):
                set_text_in_shape(shape, "29€/mes · Básico")
            elif shape.left < Inches(9):
                set_text_in_shape(shape, "49€/mes · Intermedio")
            else:
                set_text_in_shape(shape, "79€/mes · Completo")


# ══════════════════════════════════════════════
# SLIDE 11 — 4 SCREENSHOTS (layout: '10_Title Slide')
# 4 image PHs: idx=10, 11, 12, 13
# ══════════════════════════════════════════════
print("Slide 11: Módulos")
s = prs.slides[10]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/calendario.png")
insert_screenshot_cropped(s, 11, f"{SCREENSHOTS}/clases.png")
insert_screenshot_cropped(s, 12, f"{SCREENSHOTS}/recuperaciones.png")
insert_screenshot_cropped(s, 13, f"{SCREENSHOTS}/notificaciones.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "Success Stories" in txt:
            set_text_in_shape(shape, "Módulos Principales")
        elif "AI Case Studies" in txt:
            set_text_in_shape(shape, "Funcionalidades")
        elif "AI Business Strategy leverages" in txt:
            set_text_in_shape(shape, "Cuatro módulos clave: Calendario semanal con vista interactiva, Gestión de clases con filtros avanzados, Sistema de recuperaciones automáticas y Centro de notificaciones para comunicación con alumnos.")


# ══════════════════════════════════════════════
# SLIDE 12 — SCREENSHOT CENTRADO (layout: '1_Custom Layout')
# Center image PH idx=10 (4.68"x5.76")
# ══════════════════════════════════════════════
print("Slide 12: Recuperaciones")
s = prs.slides[11]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/recuperaciones.png")


# ══════════════════════════════════════════════
# SLIDE 13 — BREAK/QUOTE (layout: '12_Title Slide')
# Image PH idx=10 (center), text around
# ══════════════════════════════════════════════
print("Slide 13: Modelo SaaS")
s = prs.slides[12]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/planes.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "Break Slide" in txt:
            set_text_in_shape(shape, "Modelo SaaS")
        elif "Reflect" in txt:
            set_text_in_shape(shape, "3 Planes de Suscripción")
        elif "60 Minutes" in txt:
            set_text_in_shape(shape, "Starter · Club · Elite")
        elif "Keeping Up" in txt:
            set_text_in_shape(shape, "Feature Gating por Plan")
        elif "The future belongs" in txt:
            set_text_in_shape(shape, '"Cada academia accede exactamente a las funcionalidades que necesita, con la posibilidad de escalar en cualquier momento."')


# ══════════════════════════════════════════════
# SLIDE 14 — DATOS / STATS (layout: 'Title Slide')
# Stats boxes (290+, 2M+, 13M, 344+, 455+)
# ══════════════════════════════════════════════
print("Slide 14: Datos Reales")
s = prs.slides[13]
stat_map14 = {
    "290+": "8",
    "2M+": "3",
    "13M": "4",
    "344+": "3",
    "455+": "44%",
    "Step-by-Step Integration": "Datos en Tiempo Real",
    "Implementation Roadmap": "Cuadro de Clases",
    "Strategic Intelligence": "",
}
desc_done14 = False
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        for old, new in stat_map14.items():
            if txt == old:
                set_text_in_shape(shape, new)
                break
        if "Step-by-Step" in txt:
            set_text_in_shape(shape, "Datos en Tiempo Real")
        elif "Implementation Roadmap" in txt:
            set_text_in_shape(shape, "Métricas de la Plataforma")
        elif "AI Business Strategy leverages" in txt and not desc_done14:
            set_text_in_shape(shape, "8 clases activas · 3 profesores (Carlos, Laura, Sergio) · 4 recuperaciones pendientes · 3 pistas · Ocupación 44% · 14 alumnos")
            desc_done14 = True

# Fix the "Strategic Intelligence" labels → real labels
si_count = 0
for shape in s.shapes:
    if hasattr(shape, "text") and shape.text.strip() == "Strategic Intelligence":
        si_count += 1
        labels = ["Clases Activas", "Profesores", "Recuperaciones", "Pistas", "Ocupación"]
        if si_count <= len(labels):
            set_text_in_shape(shape, labels[si_count - 1])
    elif hasattr(shape, "text") and shape.text.strip() == "":
        pass  # already cleared


# ══════════════════════════════════════════════
# SLIDE 15 — SCREENSHOT LATERAL (layout: '13_Title Slide')
# Image PH idx=12 + slide number
# ══════════════════════════════════════════════
print("Slide 15: Notificaciones")
s = prs.slides[14]
insert_screenshot_cropped(s, 12, f"{SCREENSHOTS}/notificaciones.png")


# ══════════════════════════════════════════════
# SLIDE 16 — STAT + IMAGE (layout: '15_Title Slide')
# Image PH idx=11 + stat text
# ══════════════════════════════════════════════
print("Slide 16: Estadística")
s = prs.slides[15]
insert_screenshot_cropped(s, 11, f"{SCREENSHOTS}/alumnos.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "23 HZ" in txt:
            set_text_in_shape(shape, "14")
        elif "Real-Time Optimization" in txt:
            set_text_in_shape(shape, "Alumnos Registrados")


# ══════════════════════════════════════════════
# SLIDE 17 — CONTACTO/CIERRE (layout: '16_Title Slide')
# Wide image PH at bottom, contact info, quote
# ══════════════════════════════════════════════
print("Slide 17: Cierre")
s = prs.slides[16]
insert_screenshot_cropped(s, 11, f"{SCREENSHOTS}/dashboard_elite.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "loremipsum@mail" in txt:
            set_text_in_shape(shape, "info@academiapadel.com")
        elif "+0123" in txt:
            set_text_in_shape(shape, "+34 912 345 678")
        elif "www.yourwebsite" in txt:
            set_text_in_shape(shape, "app.academiapadel.com")
        elif "Harnessing" in txt:
            set_text_in_shape(shape, '"La tecnología al servicio del deporte: gestión inteligente para academias de pádel"')
        elif "Let's Connect" in txt:
            set_text_in_shape(shape, "¿Preguntas?")
        elif "Solar Energy" in txt:
            set_text_in_shape(shape, "Academia Pádel · SaaS")
        elif "AI Business Strategy leverages" in txt:
            set_text_in_shape(shape, "Plataforma desarrollada con Next.js 15, Express v5, PostgreSQL 16 y Docker. Modelo SaaS con 3 planes: Starter (29€), Club (49€) y Elite (79€).")
        elif "233+" in txt:
            set_text_in_shape(shape, "826€")


# ══════════════════════════════════════════════
# SLIDE 18 — CIERRE FINAL (layout: '1_Title Slide')
# Full bg image + quote
# ══════════════════════════════════════════════
print("Slide 18: Cierre Final")
s = prs.slides[17]
insert_screenshot_cropped(s, 10, f"{SCREENSHOTS}/dashboard_elite_full.png")
for shape in s.shapes:
    if hasattr(shape, "text"):
        txt = shape.text.strip()
        if "The future belongs" in txt:
            set_text_in_shape(shape, '"Gestión eficiente, alumnos satisfechos, academia en crecimiento."')
            # Clear the "– Unknown" second run
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Unknown" in run.text:
                            run.text = ""
        elif "Thank You" in txt:
            set_text_in_shape(shape, "Gracias")
        elif "Thank you for exploring" in txt:
            set_text_in_shape(shape, "Tu academia digital empieza aquí.")
        elif "Read More" in txt:
            set_text_in_shape(shape, "")


# ──────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\n✅ Presentación guardada: {OUTPUT}")
print(f"   {len(prs.slides)} diapositivas — estilo de plantilla preservado")
