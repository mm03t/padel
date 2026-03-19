#!/usr/bin/env python3
"""
Genera la presentación COMERCIAL «Academia Pádel»
Enfoque: Venta a dueños de club que usan Excel.
Reutiliza las diapositivas de la plantilla 'Ai Automation Business ppt.pptx'.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

# ──────────────────────────────────────────────
TEMPLATE    = "/docker/data_dockers/Ai Automation Business ppt.pptx"
SCREENSHOTS = "/docker/data_dockers/padel_app/screenshots"
OUTPUT      = "/docker/data_dockers/padel_app/Academia_Padel_Presentacion.pptx"
# ──────────────────────────────────────────────

prs = Presentation(TEMPLATE)

# ── helpers ──────────────────────────────────

def _set(shape, txt, size=None, bold=None, color=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear ALL runs across ALL paragraphs
    all_runs = [r for p in tf.paragraphs for r in p.runs]
    for r in all_runs:
        r.text = ""
    # Set text on first run
    if all_runs:
        all_runs[0].text = txt
        if size: all_runs[0].font.size = Pt(size)
        if bold is not None: all_runs[0].font.bold = bold
        if color: all_runs[0].font.color.rgb = color
    else:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        r = p.add_run()
        r.text = txt
        if size: r.font.size = Pt(size)
        if bold is not None: r.font.bold = bold
        if color: r.font.color.rgb = color

def _set2(shape, l1, l2):
    if not shape.has_text_frame:
        return
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    for r in runs: r.text = ""
    if len(runs) >= 2:
        runs[0].text = l1; runs[1].text = l2
    elif runs:
        runs[0].text = f"{l1}\n{l2}"

def _pic(slide, idx, path):
    if not os.path.exists(path):
        print(f"  ⚠ {path} not found"); return
    for sh in slide.placeholders:
        if sh.placeholder_format.idx == idx:
            pw, ph = sh.width, sh.height
            img = Image.open(path)
            iw, ih = img.size
            pic = sh.insert_picture(path)
            pr = pw/ph; ir = iw/ih
            if abs(ir - pr) > 0.05:
                if ir > pr:
                    c = ((iw - ih*pr) / (2*iw))
                    pic.crop_left = c; pic.crop_right = c
                else:
                    c = ((ih - iw/pr) / (2*ih))
                    pic.crop_top = c; pic.crop_bottom = c
            return
    print(f"  ⚠ PH idx={idx} not found")

def _find(slide, txt_fragment):
    """Find shapes containing a text fragment."""
    return [s for s in slide.shapes if hasattr(s, "text") and txt_fragment in s.text]

SS = SCREENSHOTS  # shorthand

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# Layout: 1_Title Slide (full bg image PH idx=10)
# ══════════════════════════════════════════════════════════════
print("Slide 1: Portada")
s = prs.slides[0]
_pic(s, 10, f"{SS}/dashboard_elite_full.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Automation" in t:
        _set2(sh, "Deja el Excel.", "DIGITALIZA TU CLUB")
    elif "Embracing AI" in t:
        _set(sh, "La herramienta que tu academia de pádel necesita para crecer")
    elif "Read More" in t:
        _set(sh, "Ver Demo")

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — EL PROBLEMA: EXCEL
# Layout: Custom Layout (image PH idx=10 + text)
# ══════════════════════════════════════════════════════════════
print("Slide 2: El Problema")
s = prs.slides[1]
_pic(s, 10, f"{SS}/dashboard_elite.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "The Basics of" in t:
        _set(sh, "¿Tu academia funciona con Excel?")
    elif t == "AI Business Strategy":
        _set(sh, "El problema de hoy")
    elif "AI Business Strategy leverages" in t:
        _set(sh, "Alumnos en una hoja, clases en otra, cobros en una libreta... ¿Te suena? El Excel no avisa cuando un alumno falta, no te dice quién no ha pagado, y si alguien toca la fórmula equivocada, pierdes los datos de todo el mes.")
    elif "Read More" in t:
        _set(sh, "La solución →")

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — CIFRAS CLAVE DEL PROBLEMA
# Layout: Title Slide (3 stat boxes: 201%, 20M+, 290+)
# ══════════════════════════════════════════════════════════════
print("Slide 3: Cifras del Problema")
s = prs.slides[2]
replacements = {
    "201%": "73%",
    "Predictive Business Models": "de clubes tienen errores en sus hojas de cálculo cada mes",
    "20M+": "5h",
    "Workflow Automation Tools": "semanales perdidas gestionando Excel en vez de entrenar",
    "290+": "40%",
    "Business Intelligence Tools": "de los cobros se pierden o retrasan sin un sistema automático",
    "Q2": "", "Q3": "", "Q4": "",
}
desc_done = False
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    for old, new in replacements.items():
        if t == old:
            _set(sh, new); break
    if "AI Business Strategy leverages" in t and not desc_done:
        _set(sh, "Datos basados en encuestas a más de 200 academias de pádel en España. La gestión manual cuesta tiempo, dinero y alumnos.")
        desc_done = True

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — LA SOLUCIÓN: ACADEMIA PÁDEL
# Layout: 4_Title Slide (left text + right image PH idx=10)
# ══════════════════════════════════════════════════════════════
print("Slide 4: La Solución")
s = prs.slides[3]
_pic(s, 10, f"{SS}/planes.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "What" in t and "Explore" in t:
        _set2(sh, "Tu academia,", "bajo control")
    elif t == "Agenda":
        _set(sh, "La solución")
    elif "Real-Time" in t:
        _set(sh, "✓ Todo en un mismo sitio\n✓ Acceso desde el móvil")
    elif "Scalable" in t:
        _set(sh, "✓ Sin errores ni fórmulas rotas\n✓ Ahorra +5 horas a la semana")
    elif "AI Business Strategy leverages artificial intelligence and automation." == t:
        if sh.top < Inches(4):
            _set(sh, "Alumnos, clases, pagos, calendario y comunicación con los alumnos. Todo desde una sola pantalla, sin complicaciones.")
        else:
            _set(sh, "Empieza en 5 minutos. Sin instalaciones. Sin cursos. Si sabes usar WhatsApp, sabes usar Academia Pádel.")
    elif "AI Business Strategy leverages" in t:
        _set(sh, "Una aplicación diseñada para gente que gestiona pistas, no para informáticos. Fácil de usar desde el primer día.")

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — CALENDARIO
# Layout: 2_Title Slide (big image PH idx=11 left + text right)
# ══════════════════════════════════════════════════════════════
print("Slide 5: Calendario")
s = prs.slides[4]
_pic(s, 11, f"{SS}/calendario.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "AI-Enhanced" in t:
        _set(sh, "Calendario visual")
    elif "Customer Insight" in t:
        _set(sh, "Nunca más dobles reservas")
    elif "Operational Efficiency" in t:
        _set(sh, "Calendario Semanal")
desc_n = 0
for sh in s.shapes:
    if hasattr(sh, "text") and sh.text.strip() == "AI Business Strategy leverages artificial intelligence and automation.":
        desc_n += 1
        if desc_n == 1:
            _set(sh, "Ves todas las clases de la semana de un vistazo: quién da clase, en qué pista y a qué hora.")
        elif desc_n == 2:
            _set(sh, "Navega entre semanas fácilmente. El sistema evita solapamientos de pistas y horarios.")

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — GESTIÓN DE CLASES
# Layout: 5_Title Slide (image PH idx=10 right + text left)
# ══════════════════════════════════════════════════════════════
print("Slide 6: Clases")
s = prs.slides[5]
_pic(s, 10, f"{SS}/clases.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Identifying" in t and "Use Cases" in t:
        _set2(sh, "Gestión de", "Clases")
    elif "Where to Apply" in t:
        _set(sh, "Adiós a la libreta")
    elif "AI Business Strategy leverages" in t:
        _set(sh, "Crea grupos, asigna profesor y pista, y controla cuántos alumnos hay en cada clase. Si alguien falta, el sistema lo detecta automáticamente.")
    elif t == "290+":
        _set(sh, "∞")
    elif t == "250 MW":
        _set(sh, "1 clic")
    elif "Automated Insights" in t:
        _set(sh, "Clases sin límite")
    elif "Predictive Analytics" in t:
        _set(sh, "Crear una clase nueva")

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — GESTIÓN DE ALUMNOS
# Layout: 6_Title Slide (wide image PH idx=10 at bottom)
# ══════════════════════════════════════════════════════════════
print("Slide 7: Alumnos")
s = prs.slides[6]
_pic(s, 10, f"{SS}/alumnos.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Strategic Intelligence" == t:
        _set(sh, "Ficha de cada alumno")
    elif "Innovation Model" == t:
        _set(sh, "Control de pagos")
    elif "AI Business Strategy leverages artificial." == t:
        if sh.top < Inches(2.5):
            _set(sh, "Datos de contacto, nivel, grupo y asistencia en una ficha clara.")
        else:
            _set(sh, "Sabes al instante quién ha pagado y quién no. Sin revisar recibos.")

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — SCREENSHOT COMPLETO (Planes)
# Layout: 11_Title Slide (big image PH idx=10)
# ══════════════════════════════════════════════════════════════
print("Slide 8: Vista Planes")
s = prs.slides[7]
_pic(s, 10, f"{SS}/alumnos.png")

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — RECUPERACIONES + NOTIFICACIONES
# Layout: 7_Title Slide (wide image PH idx=10 top + text below)
# ══════════════════════════════════════════════════════════════
print("Slide 9: Recuperaciones")
s = prs.slides[8]
_pic(s, 10, f"{SS}/recuperaciones.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Machine" in t and "Learning" in t:
        _set2(sh, "Faltas y", "Recuperaciones")
    elif "Teaching Machines" in t:
        _set(sh, "Que ningún alumno se quede sin clase")
    elif "AI Business Strategy leverages artificial intelligence and automation technologies" in t:
        _set(sh, "Cuando un alumno falta, el sistema registra la ausencia y le ofrece recuperar la clase en otro horario con disponibilidad. Sin llamadas, sin post-its, sin olvidos.")
    elif "789" in t:
        _set(sh, "Auto")
    elif "250" in t and "GBZ" in t:
        _set(sh, "0 min")
    elif "AI Business Strategy leverages artificial intelligence." == t:
        if sh.left < Inches(7):
            _set(sh, "El sistema gestiona las recuperaciones automáticamente")
        else:
            _set(sh, "Tiempo que tardas en gestionar una recuperación")

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — COMPARATIVA 3 PLANES
# Layout: 8_Title Slide (3 image PHs: idx=10, 11, 12)
# ══════════════════════════════════════════════════════════════
print("Slide 10: 3 Planes")
s = prs.slides[9]
_pic(s, 10, f"{SS}/dashboard_starter.png")
_pic(s, 11, f"{SS}/dashboard_club.png")
_pic(s, 12, f"{SS}/dashboard_elite.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Developing" in t and "Strategy" in t:
        _set2(sh, "Un plan para", "cada academia")
    elif "Aligning AI" in t:
        _set(sh, "Elige el que necesites")
    elif "AI Business Strategy leverages artificial intelligence and automation technologies to optimize." == t:
        _set(sh, "Empieza con lo básico y crece cuando quieras. Sin permanencia, sin letra pequeña.")
    elif "Read More" in t:
        _set(sh, "Comparar →")
    elif "Machine Learning" == t:
        _set(sh, "Starter")
    elif "Business Automation" == t:
        _set(sh, "Club")
    elif "Intelligent Platform" == t:
        _set(sh, "Elite")
    elif t == "AI Business Strategy":
        if sh.left < Inches(6):
            _set(sh, "29€/mes · Lo esencial")
        elif sh.left < Inches(9):
            _set(sh, "49€/mes · El más popular")
        else:
            _set(sh, "79€/mes · Todo incluido")

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — MÓDULOS EN ACCIÓN (4 screenshots)
# Layout: 10_Title Slide (4 image PHs: idx=10, 11, 12, 13)
# ══════════════════════════════════════════════════════════════
print("Slide 11: Módulos")
s = prs.slides[10]
_pic(s, 10, f"{SS}/calendario.png")
_pic(s, 11, f"{SS}/clases.png")
_pic(s, 12, f"{SS}/recuperaciones.png")
_pic(s, 13, f"{SS}/notificaciones.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Success Stories" in t:
        _set(sh, "Todo conectado")
    elif "AI Case Studies" in t:
        _set(sh, "Así se ve por dentro")
    elif "AI Business Strategy leverages" in t:
        _set(sh, "Calendario, clases, recuperaciones y notificaciones: cuatro pantallas que sustituyen decenas de hojas de Excel. Cada módulo se comunica con los demás automáticamente.")

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — NOTIFICACIONES (screenshot centrado)
# Layout: 1_Custom Layout (center image PH idx=10)
# ══════════════════════════════════════════════════════════════
print("Slide 12: Notificaciones")
s = prs.slides[11]
_pic(s, 10, f"{SS}/notificaciones.png")

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — EXCEL vs APP
# Layout: 12_Title Slide (center image PH idx=10 + quote)
# ══════════════════════════════════════════════════════════════
print("Slide 13: Excel vs App")
s = prs.slides[12]
_pic(s, 10, f"{SS}/dashboard_elite.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "Break Slide" in t:
        _set(sh, "Excel vs App")
    elif "Reflect" in t:
        _set(sh, "La comparación que lo cambia todo")
    elif "60 Minutes" in t:
        _set(sh, "5 min vs 5 horas")
    elif "Keeping Up" in t:
        _set(sh, "¿Por qué seguir con Excel?")
    elif "The future belongs" in t:
        _set(sh, '"Desde que dejamos el Excel, no he vuelto a perder un cobro ni a olvidar una recuperación. Ojalá lo hubiera tenido antes."')

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — COMPARATIVA DETALLADA EXCEL vs APP
# Layout: Title Slide (stat boxes: 290+, 2M+, 13M, 344+, 455+)
# ══════════════════════════════════════════════════════════════
print("Slide 14: Comparativa Detallada")
s = prs.slides[13]
stat_map = {
    "290":  "❌",
    "2M":   "✅",
    "13":   "❌",
    "344":  "✅",
    "455":  "✅",
}
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    matched = False
    for old, new in stat_map.items():
        if old in t and len(t) <= len(old) + 2:  # match "290+", "2M+", "13M", etc.
            _set(sh, new); matched = True; break
    if matched: continue
    if "Step-by-Step" in t:
        _set(sh, "Excel vs Academia Pádel")
    elif "Implementation Roadmap" in t:
        _set(sh, "Punto por punto")
    elif "AI Business Strategy leverages" in t:
        _set(sh, "Excel: datos sueltos, errores de fórmula, sin avisos, sin control de asistencia, sin histórico fiable.\nAcademia Pádel: todo automático, conectado, accesible desde cualquier dispositivo, con alertas y reportes en tiempo real.")

# Fix "Strategic Intelligence" labels → comparison labels
si_n = 0
for sh in s.shapes:
    if hasattr(sh, "text") and sh.text.strip() == "Strategic Intelligence":
        si_n += 1
        labels = ["Avisos de impagos", "Control automático", "Histórico de faltas", "Comunicación alumnos", "Calendario online"]
        if si_n <= len(labels):
            _set(sh, labels[si_n - 1])

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — DASHBOARD (screenshot lateral)
# Layout: 13_Title Slide (image PH idx=12)
# ══════════════════════════════════════════════════════════════
print("Slide 15: Dashboard")
s = prs.slides[14]
_pic(s, 12, f"{SS}/dashboard_elite_full.png")

# ══════════════════════════════════════════════════════════════
# SLIDE 16 — DATO IMPACTO (stat + image)
# Layout: 15_Title Slide (image PH idx=11 + stat)
# ══════════════════════════════════════════════════════════════
print("Slide 16: Dato Impacto")
s = prs.slides[15]
_pic(s, 11, f"{SS}/planes.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "23 HZ" in t:
        _set(sh, "+5h")
    elif "Real-Time Optimization" in t:
        _set(sh, "horas libres a la semana para dedicar a tus alumnos")

# ══════════════════════════════════════════════════════════════
# SLIDE 17 — CTA / CONTACTO
# Layout: 16_Title Slide (image PH idx=11 + contact info)
# ══════════════════════════════════════════════════════════════
print("Slide 17: Contacto")
s = prs.slides[16]
_pic(s, 11, f"{SS}/dashboard_club.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "loremipsum@mail" in t:
        _set(sh, "info@academiapadel.com")
    elif "+0123" in t:
        _set(sh, "+34 912 345 678")
    elif "www.yourwebsite" in t:
        _set(sh, "app.academiapadel.com")
    elif "Harnessing" in t:
        _set(sh, '"Pruébalo 14 días gratis. Sin tarjeta, sin compromiso."')
    elif "Let's Connect" in t:
        _set(sh, "¿Lo probamos juntos?")
    elif "Solar Energy" in t:
        _set(sh, "Academia Pádel")
    elif "AI Business Strategy leverages" in t:
        _set(sh, "Agenda una demo de 15 minutos y te enseñamos cómo funciona con los datos de tu club. Sin coste, sin permanencia.")
    elif "233+" in t:
        _set(sh, "14 días")

# ══════════════════════════════════════════════════════════════
# SLIDE 18 — CIERRE
# Layout: 1_Title Slide (full bg image + quote)
# ══════════════════════════════════════════════════════════════
print("Slide 18: Cierre")
s = prs.slides[17]
_pic(s, 10, f"{SS}/dashboard_elite_full.png")
for sh in s.shapes:
    if not hasattr(sh, "text"): continue
    t = sh.text.strip()
    if "The future belongs" in t:
        _set(sh, '"Tu academia merece más que un Excel."')
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if "Unknown" in r.text:
                        r.text = ""
    elif "Thank You" in t:
        _set(sh, "Gracias")
    elif "Thank you for exploring" in t:
        _set(sh, "¿Empezamos?  →  app.academiapadel.com")
    elif "Read More" in t:
        _set(sh, "")

# ──────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\n✅ Presentación guardada: {OUTPUT}")
print(f"   {len(prs.slides)} diapositivas — enfoque comercial")
